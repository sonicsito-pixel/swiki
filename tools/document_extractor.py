import os
import sys
import json
import re

# user site-packages 경로 명시적 추가 (호환성 보장)
sys.path.append(r"C:\Users\user\AppData\Roaming\Python\Python314\site-packages")

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pypdf
except ImportError:
    pypdf = None


def parse_txt(filepath):
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"Error reading TXT: {e}"


def parse_html(filepath):
    if not BeautifulSoup:
        return "BeautifulSoup4 not available."
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()
        soup = BeautifulSoup(html_content, "html.parser")
        
        # 스크립트 및 스타일 제거
        for script in soup(["script", "style"]):
            script.decompose()
            
        # 텍스트만 깔끔하게 정제
        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception as e:
        return f"Error parsing HTML: {e}"


def parse_docx(filepath):
    if not docx:
        return "python-docx not available."
    try:
        doc = docx.Document(filepath)
        output = []
        for p in doc.paragraphs:
            if p.text.strip():
                output.append(p.text.strip())
                
        for table in doc.tables:
            output.append("\n[Table]")
            for row in table.rows:
                row_vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                output.append(" | ".join(row_vals))
            output.append("[/Table]\n")
            
        return "\n".join(output)
    except Exception as e:
        return f"Error parsing DOCX: {e}"


def parse_pptx(filepath):
    if not Presentation:
        return "python-pptx not available."
    try:
        prs = Presentation(filepath)
        output = []
        for i, slide in enumerate(prs.slides):
            output.append(f"\n--- Slide {i+1} ---")
            
            # 슬라이드 내 모든 셰이프 탐색
            for shape in slide.shapes:
                # 1. 텍스트 박스인 경우
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            output.append(text)
                            
                # 2. 표(Table)인 경우
                elif shape.has_table:
                    table = shape.table
                    output.append("\n[Table]")
                    for row in table.rows:
                        row_vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        output.append(" | ".join(row_vals))
                    output.append("[/Table]\n")
                    
        return "\n".join(output)
    except Exception as e:
        return f"Error parsing PPTX: {e}"


def parse_xlsx(filepath):
    if not openpyxl:
        return "openpyxl not available."
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        output = []
        for sheetname in wb.sheetnames:
            output.append(f"\n=========================================")
            output.append(f" Sheet: {sheetname}")
            output.append(f"=========================================")
            sheet = wb[sheetname]
            
            # 빈 시트 체크
            if sheet.max_row == 1 and sheet.max_column == 1 and sheet.cell(row=1, column=1).value is None:
                continue
                
            for r in range(1, sheet.max_row + 1):
                row_vals = [sheet.cell(row=r, column=c).value for c in range(1, sheet.max_column + 1)]
                if any(v is not None for v in row_vals):
                    row_str = " | ".join([str(v).replace("\n", " ").strip() if v is not None else "" for v in row_vals])
                    output.append(f"R{r:02d}: {row_str}")
        return "\n".join(output)
    except Exception as e:
        return f"Error parsing XLSX: {e}"


def parse_pdf(filepath):
    if not pypdf:
        return "pypdf not available."
    try:
        reader = pypdf.PdfReader(filepath)
        output = []
        for i, page in enumerate(reader.pages):
            output.append(f"\n--- Page {i+1} ---")
            text = page.extract_text()
            if text:
                output.append(text.strip())
        return "\n".join(output)
    except Exception as e:
        return f"Error parsing PDF: {e}"


def extract_document(filepath):
    _, ext = os.path.splitext(filepath)
    ext = ext.lower()
    
    if ext in [".txt", ".log"]:
        return parse_txt(filepath)
    elif ext in [".html", ".htm"]:
        return parse_html(filepath)
    elif ext in [".docx", ".doc"]:
        return parse_docx(filepath)
    elif ext in [".pptx", ".ppt"]:
        return parse_pptx(filepath)
    elif ext in [".xlsx", ".xls", ".xlsm"]:
        return parse_xlsx(filepath)
    elif ext in [".pdf"]:
        return parse_pdf(filepath)
    else:
        return f"Unsupported file extension: {ext}"


def main():
    if len(sys.argv) < 2:
        print("Usage: python document_extractor.py <file_path> [output_path]")
        sys.exit(1)
        
    src_file = sys.argv[1]
    
    if not os.path.exists(src_file):
        print(f"File not found: {src_file}")
        sys.exit(1)
        
    content = extract_document(src_file)
    
    # 출력 경로가 지정되어 있으면 저장, 없으면 stdout 출력
    if len(sys.argv) >= 3:
        dst_file = sys.argv[2]
        os.makedirs(os.path.dirname(os.path.abspath(dst_file)), exist_ok=True)
        with open(dst_file, "w", encoding="utf-8") as f:
            f.write(content)
        print(f"Success. Extracted content saved to {dst_file}")
    else:
        print(content)


if __name__ == "__main__":
    main()
